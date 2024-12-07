# from fastapi import FastAPI, File, UploadFile
# from fastapi.middleware.cors import CORSMiddleware
# import magic

# app = FastAPI()

# # Enable CORS for all origins (you can restrict this in production)
# app.add_middleware(
#     CORSMiddleware,
#     allow_origins=["*"],  # Replace "*" with the frontend domain, e.g., ["http://localhost:5500"]
#     allow_credentials=True,
#     allow_methods=["*"],
#     allow_headers=["*"],
# )

# # Initialize magic
# magic_instance = magic.Magic()

# # File upload and identification endpoint
# @app.post("/identify")
# async def identify_file(file: UploadFile = File(...)):
#     try:
#         # Read the file
#         file_content = await file.read()

#         # Identify the MIME type using libmagic
#         mime_type = magic.from_buffer(file_content, mime=True)
#         file_type = magic_instance.from_buffer(file_content)
        
#         # Return result as JSON
#         return {"file_type": file_type, "mime_type": mime_type}
#     except Exception as e:
#         return {"error": str(e)}




import os
import json
import base64
import magic
import cv2
from PIL import Image
from werkzeug.utils import secure_filename
from flask import Flask, request, jsonify

from flask_cors import CORS

app = Flask(__name__)
CORS(app)  # This will allow all domains by default

# Directory to save uploaded files temporarily
UPLOAD_FOLDER = './uploads'
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
class MasterApp:
    def __init__(self):
        self.magic = magic.Magic()

    def process_file(self, file_path):
        file_format = self.identify_file_format(file_path)
        content = {}
        if "text" in file_format.lower():
            content = self.process_text(file_path)
        elif "pdf" in file_format.lower():
            content = self.process_pdf(file_path)
        elif "word" in file_format.lower():
            content = self.process_docx(file_path)
        elif "powerpoint" in file_format.lower():
            content = self.process_pptx(file_path)
        else:
            return {"error": "Unsupported file type."}
        output_path = os.path.splitext(file_path)[0] + ".json"
        self.save_to_json(content, output_path)
        return content

    def identify_file_format(self, file_path):
        return self.magic.from_file(file_path)

    def save_to_json(self, content, output_path):
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(content, f, ensure_ascii=False, indent=4)

    def process_text(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            text_content = file.read()

            return {
                "page_1": {
                    "text": {"content": text_content},
                    "recognized_text": {},
                    "images": [],
                    "tables": {}
                }
            }

    def process_pdf(self, file_path):
        import fitz  # PyMuPDF
        doc = fitz.open(file_path)
        content = {}
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            images = []

            for img_index, img in enumerate(page.get_images(full=True)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                img_bytes = base_image["image"]
                image_base64 = base64.b64encode(img_bytes).decode('utf-8')
                images.append({"base64": image_base64})

            content[f"page_{page_num + 1}"] = {
                "text": {"content": text},
                "recognized_text": {},
                "images": images,
                "tables": {}
            }
        doc.close()
        return content

    def process_docx(self, file_path):
        from docx import Document
        doc = Document(file_path)
        content = {}
        page_counter = 1

        for para in doc.paragraphs:
            content[f"page_{page_counter}"] = {
                "text": {"content": para.text.strip()},
                "recognized_text": {},
                "images": [],
                "tables": {}
            }
            page_counter += 1

        return content

    def process_pptx(self, file_path):
        from pptx import Presentation
        presentation = Presentation(file_path)
        content = {}
        page_counter = 1

        for slide in presentation.slides:
            text_content = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text_content.append(para.text.strip())
            content[f"page_{page_counter}"] = {
                "text": {"content": " ".join(text_content)},
                "recognized_text": {},
                "images": [],
                "tables": {}
            }
            page_counter += 1

        return content

# API Routes
@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file provided"}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    filename = secure_filename(file.filename)
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(file_path)

    content = app.master_app.process_file(file_path)
    return jsonify(content)

if __name__ == "__main__":
    app.master_app = MasterApp()
    app.run(host='0.0.0.0', port=5000)
